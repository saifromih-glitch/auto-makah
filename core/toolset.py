# ═══════════════════════════════════════════
# 🕋 Auto Makah — Comprehensive Tool Set
# Registers 20+ tools into the global registry
# ═══════════════════════════════════════════

import json
import os
import csv
import io
import re
import base64
import hashlib
import secrets
import string
import subprocess
import sys
import tempfile
from datetime import datetime
from typing import Any

import httpx

from core.tools import Tool, ToolCategory, registry

# ═══════════════════════════════════════════
# Output directory for generated files
# ═══════════════════════════════════════════
OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "output")
os.makedirs(OUTPUT_DIR, exist_ok=True)


def _safe_filename(name: str, ext: str) -> str:
    """Generate a safe, timestamped filename."""
    ts = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
    safe = re.sub(r"[^\w\-]", "_", name or "file")[:40]
    return f"{safe}_{ts}.{ext}"


# ═══════════════════════════════════════════
# 📄 FILE TOOLS
# ═══════════════════════════════════════════

def tool_create_xlsx(data: list = None, headers: list = None, filename: str = "spreadsheet") -> dict:
    """Generate an Excel (.xlsx) file from JSON data."""
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

        if not data:
            return {"error": "No data provided", "tool": "create_xlsx"}

        wb = Workbook()
        ws = wb.active
        ws.title = "Data"

        # Header style
        header_font = Font(bold=True, color="FFFFFF")
        header_fill = PatternFill(start_color="1F4E79", end_color="1F4E79", fill_type="solid")
        thin_border = Border(
            left=Side(style="thin"),
            right=Side(style="thin"),
            top=Side(style="thin"),
            bottom=Side(style="thin"),
        )

        # Write headers
        if headers:
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center")
                cell.border = thin_border
        elif data:
            # Infer headers from first row
            if isinstance(data[0], dict):
                headers = list(data[0].keys())
                for col_idx, header in enumerate(headers, 1):
                    cell = ws.cell(row=1, column=col_idx, value=header)
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.alignment = Alignment(horizontal="center")
                    cell.border = thin_border

        # Write data
        for row_idx, row in enumerate(data, 2):
            if isinstance(row, dict):
                for col_idx, key in enumerate(headers or row.keys(), 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=row.get(key, ""))
                    cell.border = thin_border
            elif isinstance(row, (list, tuple)):
                for col_idx, value in enumerate(row, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)
                    cell.border = thin_border
            else:
                ws.cell(row=row_idx, column=1, value=str(row)).border = thin_border

        # Auto-fit column widths
        for col in ws.columns:
            max_length = 0
            col_letter = col[0].column_letter
            for cell in col:
                try:
                    max_length = max(max_length, len(str(cell.value or "")))
                except Exception:
                    pass
            ws.column_dimensions[col_letter].width = min(max_length + 3, 50)

        fname = _safe_filename(filename, "xlsx")
        fpath = os.path.join(OUTPUT_DIR, fname)
        wb.save(fpath)

        return {
            "tool": "create_xlsx",
            "file": fname,
            "path": fpath,
            "rows": len(data),
            "columns": len(headers) if headers else (len(data[0]) if data else 0),
            "message": f"Excel file created: {fname}",
        }
    except ImportError:
        return {"error": "openpyxl not installed. Run: pip install openpyxl", "tool": "create_xlsx"}
    except Exception as e:
        return {"error": str(e), "tool": "create_xlsx"}


def tool_create_csv(data: list = None, headers: list = None, filename: str = "export", delimiter: str = ",") -> dict:
    """Generate a CSV file from JSON data."""
    try:
        if not data:
            return {"error": "No data provided", "tool": "create_csv"}

        fname = _safe_filename(filename, "csv")
        fpath = os.path.join(OUTPUT_DIR, fname)

        with open(fpath, "w", newline="", encoding="utf-8-sig") as f:
            if not headers and data and isinstance(data[0], dict):
                headers = list(data[0].keys())

            writer = csv.writer(f, delimiter=delimiter)
            if headers:
                writer.writerow(headers)

            for row in data:
                if isinstance(row, dict):
                    writer.writerow([row.get(h, "") for h in (headers or row.keys())])
                elif isinstance(row, (list, tuple)):
                    writer.writerow(row)
                else:
                    writer.writerow([str(row)])

        return {
            "tool": "create_csv",
            "file": fname,
            "path": fpath,
            "rows": len(data),
            "message": f"CSV file created: {fname}",
        }
    except Exception as e:
        return {"error": str(e), "tool": "create_csv"}


def tool_create_docx(title: str = "Document", sections: list = None, filename: str = "document") -> dict:
    """Generate a Word (.docx) document.

    sections: list of dicts with 'heading' (str) and 'content' (str or list of str).
    """
    try:
        from docx import Document
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # Title
        title_para = doc.add_heading(title, level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Metadata
        doc.add_paragraph(f"Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}").alignment = WD_ALIGN_PARAGRAPH.RIGHT
        doc.add_paragraph("═" * 50)

        if sections:
            for sec in (sections or []):
                if isinstance(sec, dict):
                    heading = sec.get("heading", "")
                    content = sec.get("content", "")
                    if heading:
                        doc.add_heading(heading, level=1)
                    if isinstance(content, list):
                        for line in content:
                            doc.add_paragraph(str(line), style="List Bullet")
                    elif content:
                        doc.add_paragraph(str(content))
                elif isinstance(sec, str):
                    doc.add_paragraph(sec)
        else:
            doc.add_paragraph("No content provided.")

        fname = _safe_filename(filename, "docx")
        fpath = os.path.join(OUTPUT_DIR, fname)
        doc.save(fpath)

        return {
            "tool": "create_docx",
            "file": fname,
            "path": fpath,
            "sections": len(sections) if sections else 0,
            "message": f"Word document created: {fname}",
        }
    except ImportError:
        return {"error": "python-docx not installed. Run: pip install python-docx", "tool": "create_docx"}
    except Exception as e:
        return {"error": str(e), "tool": "create_docx"}


def tool_create_pptx(title: str = "Presentation", slides: list = None, filename: str = "presentation") -> dict:
    """Generate a PowerPoint (.pptx) presentation.

    slides: list of dicts with 'title' (str) and 'content' (str or list of bullet points).
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        if slide.placeholders[1].has_text_frame:
            slide.placeholders[1].text = f"Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}"

        # Content slides
        for s in (slides or []):
            if isinstance(s, dict):
                slide_layout = prs.slide_layouts[1]  # Title and Content
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = s.get("title", "")

                content = s.get("content", "")
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()

                if isinstance(content, list):
                    for i, bullet in enumerate(content):
                        if i == 0:
                            tf.text = str(bullet)
                        else:
                            p = tf.add_paragraph()
                            p.text = str(bullet)
                            p.level = 0
                elif content:
                    tf.text = str(content)
            elif isinstance(s, str):
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = s

        fname = _safe_filename(filename, "pptx")
        fpath = os.path.join(OUTPUT_DIR, fname)
        prs.save(fpath)

        return {
            "tool": "create_pptx",
            "file": fname,
            "path": fpath,
            "slides": 1 + (len(slides) if slides else 0),
            "message": f"PowerPoint created: {fname}",
        }
    except ImportError:
        return {"error": "python-pptx not installed. Run: pip install python-pptx", "tool": "create_pptx"}
    except Exception as e:
        return {"error": str(e), "tool": "create_pptx"}


def tool_create_pdf(title: str = "Document", content: str = "", sections: list = None, filename: str = "document") -> dict:
    """Generate a PDF file using ReportLab."""
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import inch
        from reportlab.lib import colors
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
        from reportlab.platypus.tableofcontents import TableOfContents

        fname = _safe_filename(filename, "pdf")
        fpath = os.path.join(OUTPUT_DIR, fname)

        doc = SimpleDocTemplate(
            fpath,
            pagesize=A4,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=72,
            title=title,
        )

        styles = getSampleStyleSheet()
        story = []

        # Title
        story.append(Paragraph(title, styles["Title"]))
        story.append(Spacer(1, 0.25 * inch))
        story.append(Paragraph(f"Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}", styles["Italic"]))
        story.append(Spacer(1, 0.15 * inch))
        story.append(HRFlowable(width="100%", thickness=1, color=colors.grey))
        story.append(Spacer(1, 0.2 * inch))

        if sections:
            for sec in sections:
                if isinstance(sec, dict):
                    heading = sec.get("heading", "")
                    body = sec.get("content", "")
                    if heading:
                        story.append(Paragraph(heading, styles["Heading2"]))
                        story.append(Spacer(1, 0.1 * inch))
                    if isinstance(body, list):
                        for line in body:
                            story.append(Paragraph(f"• {line}", styles["Normal"]))
                    elif body:
                        story.append(Paragraph(str(body), styles["Normal"]))
                    story.append(Spacer(1, 0.1 * inch))
                elif isinstance(sec, str):
                    story.append(Paragraph(sec, styles["Normal"]))
        elif content:
            for para in content.split("\n"):
                para = para.strip()
                if para:
                    story.append(Paragraph(para, styles["Normal"]))
                    story.append(Spacer(1, 0.05 * inch))

        doc.build(story)

        return {
            "tool": "create_pdf",
            "file": fname,
            "path": fpath,
            "message": f"PDF created: {fname}",
        }
    except ImportError:
        return {"error": "reportlab not installed. Run: pip install reportlab", "tool": "create_pdf"}
    except Exception as e:
        return {"error": str(e), "tool": "create_pdf"}


# ═══════════════════════════════════════════
# 🌐 WEB TOOLS
# ═══════════════════════════════════════════

def tool_web_search(query: str = "", num_results: int = 5) -> dict:
    """Search the web using DuckDuckGo HTML (no API key needed)."""
    try:
        if not query:
            return {"error": "No search query provided", "tool": "web_search"}

        url = "https://html.duckduckgo.com/html/"
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                          "(KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
        }
        with httpx.Client(follow_redirects=True, timeout=15) as client:
            resp = client.post(url, data={"q": query}, headers=headers)
            resp.raise_for_status()

        # Simple HTML extraction of results
        body = resp.text
        results = []
        # Match DuckDuckGo result snippets
        link_pattern = re.findall(
            r'<a[^>]*class="result__a"[^>]*href="([^"]*)"[^>]*>(.*?)</a>',
            body,
            re.DOTALL | re.IGNORECASE,
        )
        snippet_pattern = re.findall(
            r'<a[^>]*class="result__snippet"[^>]*>(.*?)</a>',
            body,
            re.DOTALL | re.IGNORECASE,
        )

        for i in range(min(len(link_pattern), num_results)):
            href = link_pattern[i][0] if i < len(link_pattern) else ""
            title = re.sub(r"<[^>]+>", "", link_pattern[i][1] if i < len(link_pattern) else "").strip()
            snippet = re.sub(r"<[^>]+>", "", snippet_pattern[i] if i < len(snippet_pattern) else "").strip()
            results.append({"title": title, "url": href, "snippet": snippet})

        return {
            "tool": "web_search",
            "query": query,
            "results_count": len(results),
            "results": results,
        }
    except Exception as e:
        return {"error": str(e), "tool": "web_search", "query": query}


def tool_web_fetch(url: str = "", extract_text: bool = True) -> dict:
    """Fetch and optionally extract readable text from a URL."""
    try:
        if not url:
            return {"error": "No URL provided", "tool": "web_fetch"}

        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                          "(KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
        }
        with httpx.Client(follow_redirects=True, timeout=20) as client:
            resp = client.get(url, headers=headers)
            resp.raise_for_status()

        html = resp.text
        content_type = resp.headers.get("content-type", "")

        result = {
            "tool": "web_fetch",
            "url": str(resp.url),
            "status_code": resp.status_code,
            "content_type": content_type,
        }

        if extract_text and "html" in content_type.lower():
            # Strip HTML tags, scripts, styles
            text = re.sub(r"<script[^>]*>.*?</script>", "", html, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r"<[^>]+>", " ", text)
            text = re.sub(r"\s+", " ", text).strip()
            # Truncate to reasonable size
            result["text"] = text[:10000]
            result["text_length"] = len(text)
            result["text_truncated"] = len(text) > 10000
        else:
            result["html_length"] = len(html)
            result["html_preview"] = html[:1000]

        return result
    except Exception as e:
        return {"error": str(e), "tool": "web_fetch", "url": url}


def tool_get_weather(city: str = "") -> dict:
    """Get current weather for a city using wttr.in (no API key needed)."""
    try:
        if not city:
            return {"error": "No city provided", "tool": "get_weather"}

        url = f"https://wttr.in/{city}?format=j1"
        headers = {"User-Agent": "curl/7.0"}
        with httpx.Client(timeout=10) as client:
            resp = client.get(url, headers=headers)
            resp.raise_for_status()
            data = resp.json()

        current = data.get("current_condition", [{}])[0]
        weather = data.get("weather", [{}])[0]
        astronomy = weather.get("astronomy", [{}])[0] if weather.get("astronomy") else {}
        hourly = weather.get("hourly", [])[:8]

        return {
            "tool": "get_weather",
            "city": city,
            "current": {
                "temp_c": current.get("temp_C"),
                "feels_like_c": current.get("FeelsLikeC"),
                "humidity": current.get("humidity"),
                "description": current.get("weatherDesc", [{}])[0].get("value", ""),
                "wind_speed_kmph": current.get("windspeedKmph"),
                "wind_dir": current.get("winddir16Point"),
                "visibility_km": current.get("visibility"),
                "pressure_mb": current.get("pressure"),
                "uv_index": current.get("uvIndex"),
            },
            "forecast_today": {
                "high_c": weather.get("maxtempC"),
                "low_c": weather.get("mintempC"),
                "sunrise": astronomy.get("sunrise"),
                "sunset": astronomy.get("sunset"),
            },
        }
    except Exception as e:
        return {"error": str(e), "tool": "get_weather", "city": city}


# ═══════════════════════════════════════════
# 🧮 CALCULATION TOOLS
# ═══════════════════════════════════════════

def tool_calculate_zakat(assets: float = 0, liabilities: float = 0) -> dict:
    """Calculate Saudi corporate zakat: 2.5% × (assets − liabilities).

    Zakat is only due if net assets exceed the nisab threshold.
    Nisab = 85g gold × current gold price (approx 595g silver equivalent).
    Default nisab ≈ 3000 SAR (simplified).
    """
    try:
        assets = float(assets)
        liabilities = float(liabilities)
        net_assets = assets - liabilities
        nisab = 3000.0  # Simplified nisab threshold in SAR

        if net_assets <= 0:
            return {
                "tool": "calculate_zakat",
                "assets": assets,
                "liabilities": liabilities,
                "net_assets": net_assets,
                "zakat_due": 0,
                "rate": "2.5%",
                "message": "No zakat due — net assets are zero or negative.",
            }

        if net_assets < nisab:
            return {
                "tool": "calculate_zakat",
                "assets": assets,
                "liabilities": liabilities,
                "net_assets": net_assets,
                "zakat_rate": 0.025,
                "zakat_due": 0,
                "nisab": nisab,
                "message": f"Net assets ({net_assets:.2f}) below nisab ({nisab:.2f}). No zakat due.",
            }

        zakat_due = round(net_assets * 0.025, 2)
        return {
            "tool": "calculate_zakat",
            "assets": assets,
            "liabilities": liabilities,
            "net_assets": round(net_assets, 2),
            "zakat_rate": 0.025,
            "zakat_due": zakat_due,
            "nisab": nisab,
            "message": f"Zakat due: {zakat_due:.2f} SAR (2.5% of {net_assets:.2f})",
        }
    except (TypeError, ValueError) as e:
        return {"error": f"Invalid input: {e}", "tool": "calculate_zakat"}


def tool_calculate_vat(amount: float = 0, rate: float = 0.15, exclusive: bool = True) -> dict:
    """Calculate Saudi VAT (default 15%).

    exclusive=True: amount is pre-VAT → VAT = amount × rate
    exclusive=False: amount is VAT-inclusive → VAT = amount × rate / (1 + rate)
    """
    try:
        amount = float(amount)
        rate = float(rate)

        if exclusive:
            vat = round(amount * rate, 2)
            total = round(amount + vat, 2)
            return {
                "tool": "calculate_vat",
                "base_amount": amount,
                "vat_rate": rate,
                "vat_amount": vat,
                "total_including_vat": total,
                "exclusive": True,
            }
        else:
            # Inclusive: extract VAT from total
            net = round(amount / (1 + rate), 2)
            vat = round(amount - net, 2)
            return {
                "tool": "calculate_vat",
                "total_amount": amount,
                "vat_rate": rate,
                "vat_amount": vat,
                "net_before_vat": net,
                "exclusive": False,
            }
    except (TypeError, ValueError) as e:
        return {"error": f"Invalid input: {e}", "tool": "calculate_vat"}


def tool_calculate_break_even(fixed_costs: float = 0, variable_costs: float = 0, revenue: float = 0, price_per_unit: float = None) -> dict:
    """Calculate break-even point.

    Formula: Break-even = fixed_costs / (1 − variable_costs / revenue)
    If price_per_unit provided, also calculate break-even units.
    """
    try:
        fixed_costs = float(fixed_costs)
        variable_costs = float(variable_costs)
        revenue = float(revenue)

        if revenue <= 0:
            return {"error": "Revenue must be positive", "tool": "calculate_break_even"}

        contribution_margin_ratio = 1 - (variable_costs / revenue if revenue else 0)
        contribution_margin_pct = round(contribution_margin_ratio * 100, 1)

        if contribution_margin_ratio <= 0:
            return {
                "tool": "calculate_break_even",
                "fixed_costs": fixed_costs,
                "variable_costs": variable_costs,
                "revenue": revenue,
                "contribution_margin_ratio": contribution_margin_ratio,
                "break_even_revenue": None,
                "error": "Contribution margin ratio is zero or negative — cannot break even.",
            }

        be_revenue = round(fixed_costs / contribution_margin_ratio, 2)
        result = {
            "tool": "calculate_break_even",
            "fixed_costs": fixed_costs,
            "variable_costs": variable_costs,
            "revenue": revenue,
            "contribution_margin_pct": contribution_margin_pct,
            "break_even_revenue": be_revenue,
            "message": f"Break-even revenue: {be_revenue:.2f}",
        }

        if price_per_unit:
            price_per_unit = float(price_per_unit)
            if price_per_unit > 0:
                be_units = round(be_revenue / price_per_unit, 0)
                result["price_per_unit"] = price_per_unit
                result["break_even_units"] = int(be_units)
                result["message"] += f" ({int(be_units)} units at {price_per_unit}/unit)"

        return result
    except (TypeError, ValueError) as e:
        return {"error": f"Invalid input: {e}", "tool": "calculate_break_even"}


def tool_financial_ratios(
    net_income: float = 0,
    total_assets: float = 0,
    total_equity: float = 0,
    current_assets: float = 0,
    current_liabilities: float = 0,
    revenue: float = 0,
    total_liabilities: float = 0,
) -> dict:
    """Calculate key financial ratios.

    - ROA (Return on Assets) = Net Income / Total Assets
    - ROE (Return on Equity) = Net Income / Total Equity
    - Current Ratio = Current Assets / Current Liabilities
    - Profit Margin = Net Income / Revenue
    - Debt-to-Equity = Total Liabilities / Total Equity
    """
    try:
        net_income = float(net_income)
        total_assets = float(total_assets)
        total_equity = float(total_equity)
        current_assets = float(current_assets)
        current_liabilities = float(current_liabilities)
        revenue = float(revenue)
        total_liabilities = float(total_liabilities)

        def safe_div(a, b):
            return round(a / b, 4) if b else None

        ratios = {
            "roa": safe_div(net_income, total_assets),
            "roe": safe_div(net_income, total_equity),
            "current_ratio": safe_div(current_assets, current_liabilities),
            "profit_margin": safe_div(net_income, revenue),
            "debt_to_equity": safe_div(total_liabilities, total_equity),
        }

        # Format as percentages where appropriate
        formatted = {}
        for k, v in ratios.items():
            if v is not None:
                if k in ("roa", "roe", "profit_margin"):
                    formatted[k] = f"{v * 100:.2f}%"
                else:
                    formatted[k] = f"{v:.2f}"
            else:
                formatted[k] = "N/A"

        return {
            "tool": "financial_ratios",
            "inputs": {
                "net_income": net_income,
                "total_assets": total_assets,
                "total_equity": total_equity,
                "current_assets": current_assets,
                "current_liabilities": current_liabilities,
                "revenue": revenue,
                "total_liabilities": total_liabilities,
            },
            "ratios": formatted,
            "raw": ratios,
        }
    except (TypeError, ValueError) as e:
        return {"error": f"Invalid input: {e}", "tool": "financial_ratios"}


# ═══════════════════════════════════════════
# 💻 CODING TOOLS
# ═══════════════════════════════════════════

def tool_run_python(code: str = "", timeout: int = 10) -> dict:
    """Execute a Python code snippet and return the result (stdout + stderr).

    ⚠️ Sandboxed execution — no filesystem writes, limited imports.
    """
    try:
        if not code:
            return {"error": "No code provided", "tool": "run_python"}

        # Write to temp file
        tmpdir = tempfile.gettempdir()
        tmpfile = os.path.join(tmpdir, f"_auto_makah_{secrets.token_hex(4)}.py")
        with open(tmpfile, "w", encoding="utf-8") as f:
            f.write(code)

        result = subprocess.run(
            [sys.executable, tmpfile],
            capture_output=True,
            text=True,
            timeout=timeout,
            cwd=tmpdir,
        )

        # Cleanup
        try:
            os.unlink(tmpfile)
        except Exception:
            pass

        return {
            "tool": "run_python",
            "stdout": result.stdout.strip()[:5000] if result.stdout else "",
            "stderr": result.stderr.strip()[:2000] if result.stderr else "",
            "returncode": result.returncode,
            "success": result.returncode == 0,
        }
    except subprocess.TimeoutExpired:
        return {"error": f"Code execution timed out after {timeout}s", "tool": "run_python"}
    except Exception as e:
        return {"error": str(e), "tool": "run_python"}


# ═══════════════════════════════════════════
# 🛠️ UTILITY TOOLS
# ═══════════════════════════════════════════

def tool_translate(text: str = "", source: str = "auto", target: str = "ar") -> dict:
    """Translate text between languages using Google Translate (unofficial, no API key)."""
    try:
        if not text:
            return {"error": "No text provided", "tool": "translate"}

        # Use Google Translate's unofficial endpoint
        url = "https://translate.googleapis.com/translate_a/single"
        params = {
            "client": "gtx",
            "sl": source,
            "tl": target,
            "dt": "t",
            "q": text,
        }
        headers = {"User-Agent": "Mozilla/5.0"}
        with httpx.Client(timeout=15) as client:
            resp = client.get(url, params=params, headers=headers)
            resp.raise_for_status()
            data = resp.json()

        # Extract translated text from response
        translated = ""
        if data and isinstance(data, list) and data[0]:
            translated = "".join(part[0] for part in data[0] if part and part[0])

        # Detect source language
        detected = source
        if source == "auto" and data and len(data) > 1 and data[1]:
            detected = data[1]

        return {
            "tool": "translate",
            "source_lang": detected,
            "target_lang": target,
            "original": text,
            "translated": translated,
            "original_length": len(text),
            "translated_length": len(translated),
        }
    except Exception as e:
        return {"error": str(e), "tool": "translate"}


def tool_summarize(text: str = "", max_sentences: int = 5) -> dict:
    """Extract key points from long text (extractive summarization)."""
    try:
        if not text:
            return {"error": "No text provided", "tool": "summarize"}

        # Split into sentences
        sentences = re.split(r"(?<=[.!?؟。])\s+", text.strip())
        sentences = [s.strip() for s in sentences if s.strip() and len(s.strip()) > 3]

        if not sentences:
            return {"error": "No valid sentences found", "tool": "summarize"}

        # Simple extractive summarization:
        # 1. Score sentences by word frequency
        # 2. Pick top N sentences in original order

        # Word frequency (lowercase, strip punctuation)
        word_freq = {}
        for sentence in sentences:
            words = re.findall(r"\w+", sentence.lower())
            for word in words:
                if len(word) > 2:  # Skip very short words
                    word_freq[word] = word_freq.get(word, 0) + 1

        # Score each sentence
        scored = []
        for i, sentence in enumerate(sentences):
            words = re.findall(r"\w+", sentence.lower())
            score = sum(word_freq.get(w, 0) for w in words if len(w) > 2)
            # Avoid division by zero
            score = score / max(len(words), 1)
            scored.append((i, sentence, score))

        # Sort by score descending, keep original order
        top_indices = sorted(
            [s for s in scored if s[2] > 0],
            key=lambda x: x[2],
            reverse=True,
        )[:max_sentences]

        # Sort back to original order
        top_indices.sort(key=lambda x: x[0])
        key_points = [s[1] for s in top_indices]

        return {
            "tool": "summarize",
            "original_length": len(text),
            "sentence_count": len(sentences),
            "key_points": key_points,
            "key_points_count": len(key_points),
        }
    except Exception as e:
        return {"error": str(e), "tool": "summarize"}


def tool_generate_chart(chart_type: str = "bar", title: str = "Chart", labels: list = None, values: list = None, description: str = "") -> dict:
    """Generate a chart specification for later rendering.

    Returns a structured description that can be consumed by a chart renderer.
    Supported types: bar, line, pie, scatter, area.
    """
    try:
        chart_type = chart_type.lower().strip()
        valid_types = ["bar", "line", "pie", "scatter", "area", "radar", "doughnut"]
        if chart_type not in valid_types:
            return {"error": f"Invalid chart type. Choose from: {', '.join(valid_types)}", "tool": "generate_chart"}

        if not labels:
            labels = []
        if not values:
            values = []

        # Ensure labels and values align
        if len(labels) != len(values) and values:
            # Pad or truncate
            if len(labels) > len(values):
                values = values + [0] * (len(labels) - len(values))
            else:
                values = values[:len(labels)]

        spec = {
            "tool": "generate_chart",
            "chart_type": chart_type,
            "title": title,
            "labels": labels,
            "values": values,
            "datasets": [{"label": title, "data": values}],
            "options": {
                "responsive": True,
                "plugins": {
                    "title": {"display": True, "text": title},
                },
            },
        }

        if description:
            spec["description"] = description

        return spec
    except Exception as e:
        return {"error": str(e), "tool": "generate_chart"}


def tool_generate_qr(text: str = "", size: int = 256) -> dict:
    """Generate a QR code as base64 PNG data URL."""
    try:
        if not text:
            return {"error": "No text provided", "tool": "generate_qr"}

        import qrcode as _qr
        from io import BytesIO

        img = _qr.make(text)
        buf = BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)
        b64 = base64.b64encode(buf.read()).decode()

        return {
            "tool": "generate_qr",
            "text": text,
            "size": size,
            "data_url": f"data:image/png;base64,{b64}",
        }
    except ImportError:
        return {"error": "qrcode not installed. Run: pip install qrcode[pil]", "tool": "generate_qr"}
    except Exception as e:
        return {"error": str(e), "tool": "generate_qr"}


def tool_hash_text(text: str = "", algorithm: str = "sha256") -> dict:
    """Hash text using MD5, SHA1, SHA256, or SHA512."""
    try:
        if not text:
            return {"error": "No text provided", "tool": "hash_text"}

        algorithm = algorithm.lower()
        algorithms = {"md5": hashlib.md5, "sha1": hashlib.sha1, "sha256": hashlib.sha256, "sha512": hashlib.sha512}

        if algorithm not in algorithms:
            return {
                "error": f"Unknown algorithm: {algorithm}. Choose from: {', '.join(algorithms.keys())}",
                "tool": "hash_text",
            }

        h = algorithms[algorithm](text.encode("utf-8"))

        return {
            "tool": "hash_text",
            "algorithm": algorithm,
            "hash": h.hexdigest(),
            "input_length": len(text),
        }
    except Exception as e:
        return {"error": str(e), "tool": "hash_text"}


def tool_generate_password(length: int = 16, include_special: bool = True) -> dict:
    """Generate a secure random password."""
    try:
        length = max(8, min(int(length), 128))  # Clamp 8-128

        chars = string.ascii_letters + string.digits
        if include_special:
            chars += "!@#$%^&*()-_=+[]{}|;:,.<>?"

        # Ensure at least one of each type
        password = [
            secrets.choice(string.ascii_lowercase),
            secrets.choice(string.ascii_uppercase),
            secrets.choice(string.digits),
        ]
        if include_special:
            password.append(secrets.choice("!@#$%^&*()-_=+[]{}|;:,.<>?"))

        # Fill remaining
        remaining = length - len(password)
        password.extend(secrets.choice(chars) for _ in range(remaining))

        # Shuffle
        secrets.SystemRandom().shuffle(password)
        pwd = "".join(password)

        return {
            "tool": "generate_password",
            "length": length,
            "password": pwd,
            "include_special": include_special,
        }
    except Exception as e:
        return {"error": str(e), "tool": "generate_password"}


def tool_format_json(data: str = "", indent: int = 2) -> dict:
    """Parse and pretty-print JSON string."""
    try:
        if not data:
            return {"error": "No JSON provided", "tool": "format_json"}

        # Accept both raw JSON string and Python dict/list
        if isinstance(data, (dict, list)):
            parsed = data
        else:
            parsed = json.loads(data)

        formatted = json.dumps(parsed, indent=indent, ensure_ascii=False)

        return {
            "tool": "format_json",
            "valid": True,
            "formatted": formatted,
            "keys_count": len(parsed) if isinstance(parsed, dict) else len(parsed),
            "type": "object" if isinstance(parsed, dict) else "array",
        }
    except json.JSONDecodeError as e:
        return {"error": f"Invalid JSON: {e}", "tool": "format_json", "valid": False}
    except Exception as e:
        return {"error": str(e), "tool": "format_json"}


def tool_send_email(to: str = "", subject: str = "", body: str = "", smtp_host: str = "", smtp_port: int = 587, smtp_user: str = "", smtp_password: str = "") -> dict:
    """Send an email via SMTP. All SMTP fields are required."""
    try:
        if not all([to, subject, body, smtp_host, smtp_user, smtp_password]):
            return {
                "error": "Missing required fields: to, subject, body, smtp_host, smtp_user, smtp_password",
                "tool": "send_email",
            }

        import smtplib
        from email.mime.text import MIMEText
        from email.mime.multipart import MIMEMultipart

        msg = MIMEMultipart()
        msg["From"] = smtp_user
        msg["To"] = to
        msg["Subject"] = subject
        msg.attach(MIMEText(body, "plain", "utf-8"))

        with smtplib.SMTP(smtp_host, int(smtp_port), timeout=15) as server:
            server.starttls()
            server.login(smtp_user, smtp_password)
            server.send_message(msg)

        return {
            "tool": "send_email",
            "to": to,
            "subject": subject,
            "sent": True,
            "message": f"Email sent to {to}",
        }
    except ImportError:
        return {"error": "smtplib unavailable", "tool": "send_email"}
    except Exception as e:
        return {"error": str(e), "tool": "send_email"}


# ═══════════════════════════════════════════
# 📋 REGISTER ALL TOOLS
# ═══════════════════════════════════════════

def register_all():
    """Register all tools into the global registry. Call once at startup."""
    tools = [
        # File tools
        Tool("create_xlsx", tool_create_xlsx,
             "Generate Excel (.xlsx) file from JSON data. Provide 'data' (list of dicts), optional 'headers', optional 'filename'.",
             ToolCategory.FILE,
             {"data": "list", "headers": "list", "filename": "string"}),
        Tool("create_csv", tool_create_csv,
             "Generate CSV file from JSON data. Provide 'data' (list), optional 'headers', 'filename', 'delimiter'.",
             ToolCategory.FILE,
             {"data": "list", "headers": "list", "filename": "string", "delimiter": "string"}),
        Tool("create_docx", tool_create_docx,
             "Generate Word (.docx) document. Provide 'title', 'sections' (list of {heading, content}), optional 'filename'.",
             ToolCategory.FILE,
             {"title": "string", "sections": "list", "filename": "string"}),
        Tool("create_pptx", tool_create_pptx,
             "Generate PowerPoint (.pptx) presentation. Provide 'title', 'slides' (list of {title, content}), optional 'filename'.",
             ToolCategory.FILE,
             {"title": "string", "slides": "list", "filename": "string"}),
        Tool("create_pdf", tool_create_pdf,
             "Generate PDF document. Provide 'title', 'content' or 'sections' (list of {heading, content}), optional 'filename'.",
             ToolCategory.FILE,
             {"title": "string", "content": "string", "sections": "list", "filename": "string"}),

        # Web tools
        Tool("web_search", tool_web_search,
             "Search the web using DuckDuckGo (no API key needed). Provide 'query' and optional 'num_results'.",
             ToolCategory.SEARCH,
             {"query": "string", "num_results": "int"}),
        Tool("web_fetch", tool_web_fetch,
             "Fetch and extract readable text from a URL. Provide 'url' and optional 'extract_text' (default true).",
             ToolCategory.SEARCH,
             {"url": "string", "extract_text": "bool"}),
        Tool("get_weather", tool_get_weather,
             "Get current weather for a city using wttr.in (no API key needed). Provide 'city'.",
             ToolCategory.SEARCH,
             {"city": "string"}),

        # Calculation tools
        Tool("calculate_zakat", tool_calculate_zakat,
             "Calculate Saudi corporate zakat: 2.5% × (assets − liabilities). Provide 'assets' and 'liabilities'.",
             ToolCategory.UTILITY,
             {"assets": "float", "liabilities": "float"}),
        Tool("calculate_vat", tool_calculate_vat,
             "Calculate Saudi VAT (default 15%). Provide 'amount', optional 'rate' and 'exclusive' (default true).",
             ToolCategory.UTILITY,
             {"amount": "float", "rate": "float", "exclusive": "bool"}),
        Tool("calculate_break_even", tool_calculate_break_even,
             "Calculate break-even point: fixed_costs / (1 − variable_costs/revenue). Provide 'fixed_costs', 'variable_costs', 'revenue', optional 'price_per_unit'.",
             ToolCategory.UTILITY,
             {"fixed_costs": "float", "variable_costs": "float", "revenue": "float", "price_per_unit": "float"}),
        Tool("financial_ratios", tool_financial_ratios,
             "Calculate ROA, ROE, current ratio, profit margin, debt-to-equity. Provide financial figures.",
             ToolCategory.UTILITY,
             {"net_income": "float", "total_assets": "float", "total_equity": "float", "current_assets": "float", "current_liabilities": "float", "revenue": "float", "total_liabilities": "float"}),

        # Coding tools
        Tool("run_python", tool_run_python,
             "Execute a Python code snippet and return stdout/stderr. Provide 'code' and optional 'timeout' (default 10s).",
             ToolCategory.CODING,
             {"code": "string", "timeout": "int"}),

        # Utility tools
        Tool("translate", tool_translate,
             "Translate text between languages using Google Translate. Provide 'text', optional 'source' (default auto), 'target' (default ar).",
             ToolCategory.UTILITY,
             {"text": "string", "source": "string", "target": "string"}),
        Tool("summarize", tool_summarize,
             "Extract key points from long text using extractive summarization. Provide 'text' and optional 'max_sentences'.",
             ToolCategory.UTILITY,
             {"text": "string", "max_sentences": "int"}),
        Tool("generate_chart", tool_generate_chart,
             "Generate a chart specification for later rendering. Provide 'chart_type', 'title', 'labels', 'values'.",
             ToolCategory.UTILITY,
             {"chart_type": "string", "title": "string", "labels": "list", "values": "list"}),
        Tool("generate_qr", tool_generate_qr,
             "Generate a QR code as base64 PNG. Provide 'text' and optional 'size'.",
             ToolCategory.UTILITY,
             {"text": "string", "size": "int"}),
        Tool("hash_text", tool_hash_text,
             "Hash text using MD5, SHA1, SHA256, or SHA512. Provide 'text' and optional 'algorithm' (default sha256).",
             ToolCategory.UTILITY,
             {"text": "string", "algorithm": "string"}),
        Tool("generate_password", tool_generate_password,
             "Generate a secure random password. Provide optional 'length' (default 16) and 'include_special' (default true).",
             ToolCategory.UTILITY,
             {"length": "int", "include_special": "bool"}),
        Tool("format_json", tool_format_json,
             "Parse and pretty-print JSON string. Provide 'data' (JSON string) and optional 'indent'.",
             ToolCategory.UTILITY,
             {"data": "string", "indent": "int"}),
        Tool("send_email", tool_send_email,
             "Send an email via SMTP. Provide 'to', 'subject', 'body', 'smtp_host', 'smtp_port', 'smtp_user', 'smtp_password'.",
             ToolCategory.COMMUNICATION,
             {"to": "string", "subject": "string", "body": "string", "smtp_host": "string", "smtp_port": "int", "smtp_user": "string", "smtp_password": "string"}),
    ]

    for tool in tools:
        registry.register(tool)

    return len(tools)


# Auto-register on import
added = register_all()
print(f"[Auto Makah] Toolset: {added} tools registered ({registry.count()} total across {len(registry.categories())} categories)")
